from __future__ import annotations

import base64
import csv
import io
import re
from collections.abc import Iterable
from typing import Protocol

from app.pipeline.normalization import slug
from app.schemas import SourceFragment

try:
    import pdfplumber
except ImportError:  # pragma: no cover
    pdfplumber = None

try:
    import fitz  # PyMuPDF: рендер страниц без текстового слоя для мультимодального извлечения
except ImportError:  # pragma: no cover
    fitz = None

try:
    from docx import Document as DocxDocument
    from docx.oxml.ns import qn as docx_qn
    from docx.table import Table as DocxTable
    from docx.text.paragraph import Paragraph as DocxParagraph
except ImportError:  # pragma: no cover
    DocxDocument = None
    docx_qn = None
    DocxTable = None
    DocxParagraph = None

try:
    from openpyxl import load_workbook
except ImportError:  # pragma: no cover
    load_workbook = None

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover
    Presentation = None


class Parser(Protocol):
    name: str

    def parse(self, document_id: str, version_id: str, filename: str, content: bytes) -> list[SourceFragment]:
        ...


def normalize_text(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip().lower()


def decode_text(content: bytes) -> str:
    """Экспорты из устаревших лабораторных систем часто в Windows-1251: сначала
    строгие кодировки, замена некорректных байтов — только последний резервный
    вариант."""
    for encoding in ("utf-8-sig", "cp1251"):
        try:
            return content.decode(encoding)
        except UnicodeDecodeError:
            continue
    return content.decode("utf-8", errors="replace")


# Шрифты Symbol и Wingdings кодируются в приватной области Unicode (U+F0xx):
# извлечённый из PDF текст получает вместо «°», «≤», «→» непечатаемые коды,
# которые портят и цитату, и эмбеддинг фрагмента. Таблица восстанавливает
# стандартную раскладку Adobe Symbol (код символа + 0xF000)
_PUA_SYMBOL = {
    "\uf02d": "\u2212", "\uf044": "\u0394", "\uf061": "\u03b1", "\uf062": "\u03b2",
    "\uf064": "\u03b4", "\uf065": "\u03b5", "\uf067": "\u03b3", "\uf06c": "\u03bb",
    "\uf06d": "\u03bc", "\uf070": "\u03c0", "\uf072": "\u03c1", "\uf073": "\u03c3",
    "\uf074": "\u03c4", "\uf077": "\u03c9", "\uf0a3": "\u2264", "\uf0a5": "\u221e",
    "\uf0ac": "\u2190", "\uf0ae": "\u2192", "\uf0b0": "\u00b0", "\uf0b1": "\u00b1",
    "\uf0b3": "\u2265", "\uf0b4": "\u00d7", "\uf0b7": "\u2022", "\uf0b8": "\u00f7",
    "\uf0b9": "\u2260", "\uf0bb": "\u2248", "\uf0d6": "\u221a", "\uf0d7": "\u00b7",
    "\uf0e5": "\u03a3", "\uf0fc": "\u2713",
}
_PUA_RE = re.compile("[\ue000-\uf8ff]")


def sanitize_extracted_text(text: str) -> str:
    """Восстанавливает символы шрифтов Symbol/Wingdings из приватной области.

    Код, которого нет в таблице, удаляется: он не отображается ни в одном
    шрифте и несёт не смысл, а сбой извлечения.
    """
    if not text or not _PUA_RE.search(text):
        return text
    return _PUA_RE.sub(lambda match: _PUA_SYMBOL.get(match.group(0), ""), text)


def split_text_blocks(text: str, max_chars: int = 3500) -> list[str]:
    """Разбивает текст на блоки по пустым строкам; блок длиннее max_chars
    дополнительно разбивается по границам предложений. Предложение длиннее
    лимита остаётся целым блоком — разрыв внутри фразы намеренно не выполняется."""
    raw_blocks = [block.strip() for block in re.split(r"\n\s*\n", text) if block.strip()]
    if not raw_blocks and text.strip():
        raw_blocks = [text.strip()]

    blocks: list[str] = []
    for block in raw_blocks:
        if len(block) <= max_chars:
            blocks.append(block)
            continue
        sentences = re.split(r"(?<=[.!?])\s+", block)
        current = ""
        for sentence in sentences:
            if len(current) + len(sentence) + 1 <= max_chars:
                current = f"{current} {sentence}".strip()
            else:
                if current:
                    blocks.append(current)
                current = sentence
        if current:
            blocks.append(current)
    return blocks


def merge_short_blocks(blocks: list[str], min_chars: int = 60) -> list[str]:
    """Короткие блоки (заголовок, строка оглавления, номер страницы) не остаются
    отдельными фрагментами: блок короче min_chars присоединяется к СЛЕДУЮЩЕМУ
    (заголовок — контекст абзаца) через перенос строки; если следующего нет — к
    предыдущему; пустые пропускаются. Общий помощник для PlainText/Pdf/Pptx-путей —
    тонкая обёртка над _merge_short_with_sections без разделов."""
    return [block for block, _ in _merge_short_with_sections([(block, "") for block in blocks], min_chars)]


def _merge_short_with_sections(pairs: list[tuple[str, str]], min_chars: int = 60) -> list[tuple[str, str]]:
    """Склейка коротких блоков, где каждый блок несёт свой раздел (section).
    Склеенный фрагмент получает раздел содержательного (якорного) блока —
    того, к которому присоединены предшествующие короткие заголовки/строки."""
    cleaned = [(block.strip(), section) for block, section in pairs if block and block.strip()]
    if not cleaned:
        return []
    result: list[tuple[str, str]] = []
    pending = ""  # накопленные короткие блоки, ждущие следующего содержательного
    for block, section in cleaned:
        if pending:
            block = f"{pending}\n{block}"
            pending = ""
        if len(block) < min_chars:
            pending = block
            continue
        result.append((block, section))
    if pending:
        # Оставшиеся короткие блоки без следующего присоединяются к предыдущему,
        # иначе образуют отдельный фрагмент
        if result:
            last_block, last_section = result[-1]
            result[-1] = (f"{last_block}\n{pending}", last_section)
        else:
            result.append((pending, cleaned[-1][1]))
    return result


class PlainTextParser:
    name = "plain-text"

    def parse(self, document_id: str, version_id: str, filename: str, content: bytes) -> list[SourceFragment]:
        text = decode_text(content)
        # Короткие блоки (заголовок, строка оглавления) не становятся отдельными
        # фрагментами — присоединяются к соседнему содержательному
        blocks = merge_short_blocks(split_text_blocks(text))
        return [
            SourceFragment(
                id=f"fragment-{document_id}-{index + 1}",
                document_id=document_id,
                version_id=version_id,
                page=1,
                element_type="paragraph",
                section="Uploaded text",
                text=block,
                normalized_text=normalize_text(block),
                metadata={"filename": filename, "ordinal": index + 1},
            )
            for index, block in enumerate(blocks)
        ]


class CsvParser:
    name = "csv"

    def parse(self, document_id: str, version_id: str, filename: str, content: bytes) -> list[SourceFragment]:
        text = decode_text(content)
        reader = csv.reader(io.StringIO(text))
        header_row = next(reader, None)
        if header_row is None:
            return []
        headers = _unique_headers(header_row)
        fragments: list[SourceFragment] = []
        index = 0
        for values in reader:
            if not values:
                continue
            index += 1
            # Дубли заголовков не объединяются (как это происходит в DictReader),
            # а получают суффикс; колонки сверх заголовка именуются по позиции
            row: dict[str, str] = {}
            for column, value in enumerate(values):
                key = headers[column] if column < len(headers) else f"column_{column + 1}"
                row[key] = value
            row_text = "; ".join(f"{key}={value}" for key, value in row.items() if value not in (None, ""))
            fragments.append(
                SourceFragment(
                    id=f"fragment-{document_id}-row-{index}",
                    document_id=document_id,
                    version_id=version_id,
                    page=1,
                    element_type="table_row",
                    section="Uploaded CSV",
                    text=row_text,
                    normalized_text=normalize_text(row_text),
                    metadata={"filename": filename, "row": index, "row_data": row},
                )
            )
        return fragments


class XlsxParser:
    name = "openpyxl"

    def parse(self, document_id: str, version_id: str, filename: str, content: bytes) -> list[SourceFragment]:
        if load_workbook is None:
            return BinaryPlaceholderParser(reason="openpyxl is not installed").parse(
                document_id, version_id, filename, content
            )
        try:
            workbook = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        except Exception as exc:  # pragma: no cover - depends on user files
            return BinaryPlaceholderParser(reason=f"XLSX parser failed: {exc}").parse(
                document_id, version_id, filename, content
            )

        fragments: list[SourceFragment] = []
        for sheet in workbook.worksheets:
            headers: list[str] = []
            for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                values = ["" if value is None else str(value).strip() for value in row]
                if not any(values):
                    continue
                if not headers:
                    headers = [value or f"column_{index + 1}" for index, value in enumerate(values)]
                    row_text = " | ".join(headers)
                    element_type = "xlsx_header"
                else:
                    pairs = [
                        f"{headers[index] if index < len(headers) else f'column_{index + 1}'}={value}"
                        for index, value in enumerate(values)
                        if value
                    ]
                    row_text = "; ".join(pairs)
                    element_type = "xlsx_row"
                if not row_text:
                    continue
                fragments.append(
                    SourceFragment(
                        id=f"fragment-{document_id}-xlsx-{_slug_id(sheet.title)}-{row_index}",
                        document_id=document_id,
                        version_id=version_id,
                        page=1,
                        element_type=element_type,
                        section=f"XLSX sheet {sheet.title}",
                        text=row_text,
                        normalized_text=normalize_text(row_text),
                        metadata={
                            "filename": filename,
                            "sheet": sheet.title,
                            "row": row_index,
                            "evidence_unit": True,
                            "parser": self.name,
                        },
                    )
                )

        if not fragments:
            return BinaryPlaceholderParser(reason="XLSX contains no readable cells").parse(
                document_id, version_id, filename, content
            )
        return fragments


def _page_is_mostly_image(page) -> bool:
    """Страница считается сканом, если растровые изображения покрывают
    больше половины её площади — независимо от длины текстового слоя."""
    try:
        page_area = float(page.width) * float(page.height)
        if page_area <= 0:
            return False
        covered = sum(
            max(0.0, float(image["x1"]) - float(image["x0"]))
            * max(0.0, float(image["bottom"]) - float(image["top"]))
            for image in page.images
        )
        return covered / page_area > 0.5
    except Exception:
        return False


class PdfParser:
    name = "pdfplumber"

    def parse(self, document_id: str, version_id: str, filename: str, content: bytes) -> list[SourceFragment]:
        if pdfplumber is None:
            return BinaryPlaceholderParser(reason="pdfplumber is not installed").parse(
                document_id, version_id, filename, content
            )
        fragments: list[SourceFragment] = []
        render_doc = None
        with pdfplumber.open(io.BytesIO(content)) as pdf:
            for page_index, page in enumerate(pdf.pages, start=1):
                text = sanitize_extracted_text(page.extract_text() or "").strip()
                image_b64 = None
                # Порог, а не проверка на пустоту: у скана текстовый слой часто
                # не пуст (колонтитул, OCR-штамп), но содержимое находится в
                # растровом изображении.
                # Длинный колонтитул проходит порог по тексту, поэтому вторым
                # сигналом служит растр, покрывающий большую часть страницы
                if len(text) < 40 or _page_is_mostly_image(page):
                    if fitz is not None:
                        if render_doc is None:
                            render_doc = fitz.open(stream=content, filetype="pdf")
                        try:
                            pixmap = render_doc[page_index - 1].get_pixmap(dpi=120)
                            image_b64 = base64.b64encode(pixmap.tobytes("png")).decode("ascii")
                        except Exception:
                            image_b64 = None
                    if not text:
                        text = f"PDF page {page_index} has no text layer; visual OCR adapter is required."
                blocks = merge_short_blocks(split_text_blocks(text)) or [text]
                for block_index, block in enumerate(blocks, start=1):
                    metadata = {
                        "filename": filename,
                        "page": page_index,
                        "block": block_index,
                        "evidence_unit": True,
                        "parser": self.name,
                    }
                    if image_b64 and block_index == 1:
                        metadata["image_b64"] = image_b64
                    fragments.append(
                        SourceFragment(
                            id=f"fragment-{document_id}-p{page_index}-{block_index}",
                            document_id=document_id,
                            version_id=version_id,
                            page=page_index,
                            element_type="pdf_page_text",
                            section="PDF evidence unit",
                            text=block,
                            normalized_text=normalize_text(block),
                            metadata=metadata,
                        )
                    )
        if render_doc is not None:
            render_doc.close()
        return fragments


# Подпись таблицы в отраслевых отчётах: «Таблица 1 - …», «Табл. 2 …», «Table 3 …»
_TABLE_CAPTION_RE = re.compile(r"^\s*(таблица|табл\.?|table)\s*\d", re.IGNORECASE)


def _iter_docx_body(document) -> Iterable:
    """Абзацы и таблицы в том порядке, в каком они лежат в файле.

    python-docx отдаёт document.paragraphs и document.tables отдельными
    списками, теряя их взаимное расположение; порядок восстанавливается обходом
    тела документа.
    """
    for child in document.element.body.iterchildren():
        if child.tag == docx_qn("w:p"):
            yield DocxParagraph(child, document)
        elif child.tag == docx_qn("w:tbl"):
            yield DocxTable(child, document)


class DocxParser:
    name = "python-docx"

    def parse(self, document_id: str, version_id: str, filename: str, content: bytes) -> list[SourceFragment]:
        if DocxDocument is None:
            return BinaryPlaceholderParser(reason="python-docx is not installed").parse(
                document_id, version_id, filename, content
            )
        try:
            document = DocxDocument(io.BytesIO(content))
        except Exception as exc:  # pragma: no cover - depends on user files
            return BinaryPlaceholderParser(reason=f"DOCX parser failed: {exc}").parse(
                document_id, version_id, filename, content
            )

        fragments: list[SourceFragment] = []
        # У DOCX нет страниц: адрес фрагмента — единый сквозной номер блока,
        # общий для абзацев и табличных строк (иначе адреса конфликтуют)
        block_index = 0
        table_index = 0
        section = "DOCX evidence unit"
        # Абзацы накапливаются до ближайшей таблицы или конца документа: короткие
        # (заголовок «Введение», строка оглавления) приклеиваются к следующему
        # содержательному — заголовок становится контекстом абзаца, а не
        # отдельным фрагментом. Раздел каждого блока запоминается параллельно и
        # переносится на склеенный фрагмент.
        pending: list[tuple[str, str]] = []  # (block, section)

        def flush_paragraphs() -> None:
            nonlocal block_index, pending
            for block, block_section in _merge_short_with_sections(pending):
                block_index += 1
                fragments.append(
                    SourceFragment(
                        id=f"fragment-{document_id}-docx-p{block_index}",
                        document_id=document_id,
                        version_id=version_id,
                        page=block_index,
                        element_type="docx_paragraph",
                        section=block_section,
                        text=block,
                        normalized_text=normalize_text(block),
                        metadata={
                            "filename": filename,
                            "paragraph": block_index,
                            "evidence_unit": True,
                            "parser": self.name,
                        },
                    )
                )
            pending = []

        # Тело документа обходится в порядке вёрстки: абзацы и таблицы идут
        # вперемешку, как в файле. Раздельный обход (сначала document.paragraphs,
        # затем document.tables) лишал таблицу и раздела, и подписи — по номеру
        # блока нельзя было понять, какому тексту таблица принадлежит
        for item in _iter_docx_body(document):
            if isinstance(item, DocxTable):
                # Подпись таблицы — абзац прямо перед ней; для модели это
                # единственный источник смысла колонок, кроме самой шапки
                caption = pending[-1][0] if pending else ""
                if not _TABLE_CAPTION_RE.match(caption):
                    caption = ""
                flush_paragraphs()
                table_index += 1
                for row_index, row in enumerate(item.rows, start=1):
                    cells = [sanitize_extracted_text(cell.text).strip().replace("\n", " ")
                             for cell in row.cells if cell.text.strip()]
                    row_text = " | ".join(cells)
                    if not row_text:
                        continue
                    block_index += 1
                    fragments.append(
                        SourceFragment(
                            id=f"fragment-{document_id}-docx-t{table_index}-r{row_index}",
                            document_id=document_id,
                            version_id=version_id,
                            page=block_index,
                            element_type="docx_table_row",
                            section=caption[:180] or section,
                            text=row_text,
                            normalized_text=normalize_text(row_text),
                            metadata={
                                "filename": filename,
                                "table": table_index,
                                "row": row_index,
                                "ordinal": block_index,
                                "caption": caption[:180] or None,
                                "evidence_unit": True,
                                "parser": self.name,
                            },
                        )
                    )
                continue

            text = sanitize_extracted_text(item.text).strip()
            if not text:
                continue
            style_name = (item.style.name if item.style else "") or ""
            if style_name.lower().startswith("heading"):
                section = text[:180]
            for block in split_text_blocks(text):
                pending.append((block, section))
        flush_paragraphs()

        if not fragments:
            return BinaryPlaceholderParser(reason="DOCX contains no extractable text").parse(
                document_id, version_id, filename, content
            )
        return fragments


class PptxParser:
    name = "python-pptx"

    def parse(self, document_id: str, version_id: str, filename: str, content: bytes) -> list[SourceFragment]:
        if Presentation is None:
            return BinaryPlaceholderParser(reason="python-pptx is not installed").parse(
                document_id, version_id, filename, content
            )
        try:
            presentation = Presentation(io.BytesIO(content))
        except Exception as exc:  # pragma: no cover - depends on user files
            return BinaryPlaceholderParser(reason=f"PPTX parser failed: {exc}").parse(
                document_id, version_id, filename, content
            )

        fragments: list[SourceFragment] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            # Один фрагмент на СЛАЙД: заголовок и все текстовые шейпы склеиваются
            # переносом строки (заголовок и строки оглавления не порождают
            # отдельных фрагментов). Слайд длиннее 3500 симв. разбивается, но
            # короткие строки не остаются отдельными фрагментами (merge_short_blocks).
            shape_texts = list(_slide_text_shapes(slide.shapes))
            slide_text = sanitize_extracted_text("\n".join(shape_texts)).strip()
            block_index = 0
            if slide_text:
                for block in merge_short_blocks(split_text_blocks(slide_text)) or [slide_text]:
                    block_index += 1
                    fragments.append(
                        SourceFragment(
                            id=f"fragment-{document_id}-pptx-s{slide_index}-{block_index}",
                            document_id=document_id,
                            version_id=version_id,
                            page=slide_index,
                            element_type="pptx_slide_text",
                            section=f"PPTX slide {slide_index}",
                            text=block,
                            normalized_text=normalize_text(block),
                            metadata={
                                "filename": filename,
                                "slide": slide_index,
                                "block": block_index,
                                "evidence_unit": True,
                                "parser": self.name,
                            },
                        )
                    )
            # Таблицы слайда — отдельными табличными фрагментами
            for row_index, row_text in enumerate(_slide_table_rows(slide.shapes), start=1):
                fragments.append(
                    SourceFragment(
                        id=f"fragment-{document_id}-pptx-s{slide_index}-t{row_index}",
                        document_id=document_id,
                        version_id=version_id,
                        page=slide_index,
                        element_type="pptx_table_row",
                        section=f"PPTX slide {slide_index} table",
                        text=row_text,
                        normalized_text=normalize_text(row_text),
                        metadata={
                            "filename": filename,
                            "slide": slide_index,
                            "row": row_index,
                            "evidence_unit": True,
                            "parser": self.name,
                        },
                    )
                )

        if not fragments:
            return BinaryPlaceholderParser(reason="PPTX contains no extractable text").parse(
                document_id, version_id, filename, content
            )
        return fragments


class BinaryPlaceholderParser:
    name = "binary-placeholder"

    def __init__(self, reason: str | None = None):
        self.reason = reason

    def parse(self, document_id: str, version_id: str, filename: str, content: bytes) -> list[SourceFragment]:
        text = (
            f"Файл {filename} зарегистрирован, но текст не был извлечен. "
            "Для этого формата нужен отдельный адаптер, OCR или конвертация в поддерживаемый формат."
        )
        if self.reason:
            text = f"{text} Причина: {self.reason}."
        return [
            SourceFragment(
                id=f"fragment-{document_id}-binary-1",
                document_id=document_id,
                version_id=version_id,
                page=1,
                element_type="document_placeholder",
                section="Parser adapter boundary",
                text=text,
                normalized_text=normalize_text(text),
                metadata={"filename": filename, "bytes": len(content), "parser": self.name, "reason": self.reason},
            )
        ]


# Единственный источник правды по поддерживаемым форматам:
# выбор парсера идёт через этот реестр
_PARSER_BY_EXTENSION: dict[str, type] = {
    ".csv": CsvParser,
    ".xlsx": XlsxParser,
    ".xlsm": XlsxParser,
    ".pdf": PdfParser,
    ".docx": DocxParser,
    ".docm": DocxParser,
    ".pptx": PptxParser,
    ".txt": PlainTextParser,
    ".md": PlainTextParser,
    ".json": PlainTextParser,
}


def choose_parser(filename: str) -> Parser:
    parser_class = _PARSER_BY_EXTENSION.get(extension_of(filename))
    if parser_class is not None:
        return parser_class()
    # Неизвестный формат намеренно регистрируется placeholder-фрагментом
    return BinaryPlaceholderParser(reason=f"unsupported extension {extension_of(filename) or '<none>'}")


def extension_of(filename: str) -> str:
    match = re.search(r"(\.[^.\\/]+)$", filename.lower())
    return match.group(1) if match else ""


def _unique_headers(raw: list[str]) -> list[str]:
    headers: list[str] = []
    seen: dict[str, int] = {}
    for index, value in enumerate(raw):
        name = (value or "").strip() or f"column_{index + 1}"
        count = seen.get(name, 0) + 1
        seen[name] = count
        headers.append(name if count == 1 else f"{name}_{count}")
    return headers


def _slug_id(value: str) -> str:
    # Общий slug с резервным значением: имя листа из одних спецсимволов не даёт пустой id
    return slug(value) or "sheet"


def _slide_text_shapes(shapes: Iterable) -> Iterable[str]:
    """Текст всех текстовых шейпов слайда (заголовок и прочие), без таблиц —
    они идут отдельными табличными фрагментами."""
    for shape in shapes:
        if getattr(shape, "has_text_frame", False):
            text = "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs).strip()
            if text:
                yield text
        if hasattr(shape, "shapes"):
            yield from _slide_text_shapes(shape.shapes)


def _slide_table_rows(shapes: Iterable) -> Iterable[str]:
    """Строки всех таблиц слайда (в т.ч. вложенных в группы)."""
    for shape in shapes:
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                row_text = " | ".join(cell.text.strip().replace("\n", " ") for cell in row.cells if cell.text.strip())
                if row_text:
                    yield row_text
        if hasattr(shape, "shapes"):
            yield from _slide_table_rows(shape.shapes)
